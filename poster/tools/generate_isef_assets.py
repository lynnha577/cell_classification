#!/usr/bin/env python3
"""Generate ISEF poster assets from local Allen Cell Types cache.

This script avoids external SDK dependencies by using the cached
`scripts/cell_types/ephys_features.csv` and `scripts/cell_types/cells.json`
files directly.
"""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Tuple

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests
import seaborn as sns
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestClassifier
from sklearn.impute import SimpleImputer
from sklearn.metrics import (
    ConfusionMatrixDisplay,
    PrecisionRecallDisplay,
    RocCurveDisplay,
    accuracy_score,
    average_precision_score,
    confusion_matrix,
    roc_auc_score,
)
from sklearn.model_selection import StratifiedKFold, train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.tree import DecisionTreeClassifier


plt.rcParams.update(
    {
        "figure.dpi": 180,
        "savefig.dpi": 300,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "font.size": 10,
    }
)

RANDOM_STATE = 42
CELLTYPES_BASE_URL = "https://celltypes.brain-map.org"
PROJECT_ROOT = Path(__file__).resolve().parents[2]


@dataclass
class TaskSpec:
    task_id: str
    label_col: str
    label_name: str
    top_n: int | None = None


@dataclass
class FeatureSpec:
    feature_id: str
    feature_name: str
    columns: List[str]


def safe_name(name: str) -> str:
    return re.sub(r"[^a-zA-Z0-9]+", "_", name).strip("_").lower()


def load_local_data(data_dir: Path) -> Tuple[pd.DataFrame, List[str], List[str], List[str]]:
    ephys = pd.read_csv(data_dir / "ephys_features.csv")
    cells = pd.DataFrame(json.loads((data_dir / "cells.json").read_text()))

    cells = cells.rename(
        columns={
            "specimen__id": "specimen_id",
            "donor__species": "species",
            "tag__dendrite_type": "dendrite_type",
            "structure__layer": "structure_layer_name",
            "structure_parent__acronym": "brain_area",
            "morph_thumb_path": "morph_thumb_path",
            "ephys_thumb_path": "ephys_thumb_path",
        }
    )
    cells["specimen_id"] = cells["specimen_id"].astype(int)

    merged = ephys.merge(
        cells[
            [
                "specimen_id",
                "species",
                "dendrite_type",
                "structure_layer_name",
                "brain_area",
                "morph_thumb_path",
                "ephys_thumb_path",
                "nr__average_contraction",
                "nr__average_parent_daughter_ratio",
                "nr__max_euclidean_distance",
                "nr__number_bifurcations",
                "nr__number_stems",
                "csl__normalized_depth",
            ]
        ],
        on="specimen_id",
        how="inner",
    )

    merged = merged[merged["dendrite_type"].isin(["spiny", "aspiny"])].copy()
    merged["species"] = merged["species"].fillna("Unknown")
    merged["brain_area"] = merged["brain_area"].fillna("Unknown")
    merged["structure_layer_name"] = merged["structure_layer_name"].fillna("Unknown")
    merged["species_dendrite"] = merged["species"].str.replace(" ", "", regex=False) + "_" + merged[
        "dendrite_type"
    ]
    merged["layer_group"] = merged["structure_layer_name"].replace({"6a": "6", "6b": "6", "2": "2/3", "3": "2/3"})

    drop_cols = {"id", "specimen_id", "rheobase_sweep_id", "thumbnail_sweep_id", "rheobase_sweep_number"}
    ephys_cols = [c for c in ephys.columns if c not in drop_cols]

    morph_cols = [
        "nr__average_contraction",
        "nr__average_parent_daughter_ratio",
        "nr__max_euclidean_distance",
        "nr__number_bifurcations",
        "nr__number_stems",
        "csl__normalized_depth",
    ]
    all_cols = sorted(set(ephys_cols + morph_cols))
    return merged, ephys_cols, morph_cols, all_cols


def build_pipeline(model_name: str, max_leaf_nodes: int) -> Pipeline:
    if model_name == "DecisionTree":
        model = DecisionTreeClassifier(max_leaf_nodes=max_leaf_nodes, random_state=RANDOM_STATE)
    elif model_name == "RandomForest":
        model = RandomForestClassifier(
            n_estimators=60,
            max_leaf_nodes=max_leaf_nodes,
            random_state=RANDOM_STATE,
            n_jobs=-1,
        )
    else:
        raise ValueError(f"Unknown model: {model_name}")
    return Pipeline([("imputer", SimpleImputer(strategy="median")), ("model", model)])


def evaluate_node_sweep(
    X: pd.DataFrame,
    y: pd.Series,
    node_values: np.ndarray,
    n_splits: int = 4,
) -> pd.DataFrame:
    skf = StratifiedKFold(n_splits=n_splits, shuffle=True, random_state=RANDOM_STATE)
    rows = []
    for model_name in ["DecisionTree", "RandomForest"]:
        for node in node_values:
            train_scores: List[float] = []
            test_scores: List[float] = []
            for train_idx, test_idx in skf.split(X, y):
                X_train = X.iloc[train_idx]
                X_test = X.iloc[test_idx]
                y_train = y.iloc[train_idx]
                y_test = y.iloc[test_idx]

                model = build_pipeline(model_name, int(node))
                model.fit(X_train, y_train)
                train_scores.append(model.score(X_train, y_train))
                test_scores.append(model.score(X_test, y_test))

            rows.append(
                {
                    "model": model_name,
                    "node": int(node),
                    "train_mean": float(np.mean(train_scores)),
                    "train_sem": float(np.std(train_scores, ddof=1) / np.sqrt(len(train_scores))),
                    "test_mean": float(np.mean(test_scores)),
                    "test_sem": float(np.std(test_scores, ddof=1) / np.sqrt(len(test_scores))),
                }
            )
    return pd.DataFrame(rows)


def top_classes(y: pd.Series, top_n: int | None) -> pd.Series:
    if top_n is None:
        return y
    keep = y.value_counts().head(top_n).index
    return y.where(y.isin(keep), "Other")


def plot_dataset_overview(df: pd.DataFrame, out_dir: Path) -> None:
    fig, axes = plt.subplots(1, 3, figsize=(15, 4.5))
    sns.countplot(data=df, x="species", ax=axes[0], palette="Set2")
    axes[0].set_title("Species Distribution")
    axes[0].tick_params(axis="x", rotation=20)
    axes[0].set_xlabel("")
    axes[0].set_ylabel("Cell Count")

    sns.countplot(data=df, x="dendrite_type", ax=axes[1], palette="Set1")
    axes[1].set_title("Dendrite Type Distribution")
    axes[1].set_xlabel("")
    axes[1].set_ylabel("Cell Count")

    top_areas = df["brain_area"].value_counts().head(10).index
    tmp = df[df["brain_area"].isin(top_areas)].copy()
    sns.countplot(data=tmp, y="brain_area", ax=axes[2], palette="Blues_r", order=top_areas)
    axes[2].set_title("Top 10 Brain Areas")
    axes[2].set_xlabel("Cell Count")
    axes[2].set_ylabel("")

    fig.suptitle("Allen Cell Types Dataset Snapshot", fontsize=14, fontweight="bold")
    fig.tight_layout()
    fig.savefig(out_dir / "dataset_overview.png", bbox_inches="tight")
    plt.close(fig)

    ctab = pd.crosstab(df["species"], df["dendrite_type"], normalize="index")
    fig, ax = plt.subplots(figsize=(5.8, 4.4))
    sns.heatmap(ctab, annot=True, fmt=".2f", cmap="YlOrRd", cbar_kws={"label": "Row Fraction"}, ax=ax)
    ax.set_title("Species vs Dendrite Type (Row-Normalized)")
    fig.tight_layout()
    fig.savefig(out_dir / "species_dendrite_heatmap.png", bbox_inches="tight")
    plt.close(fig)


def plot_feature_statistics(df: pd.DataFrame, ephys_cols: List[str], out_dir: Path) -> None:
    num_df = df[ephys_cols].copy()
    missing_frac = num_df.isna().mean().sort_values(ascending=False).head(20)

    fig, axes = plt.subplots(1, 2, figsize=(15, 4.6))
    missing_frac.plot(kind="bar", ax=axes[0], color="#D55E00")
    axes[0].set_title("Top Missing-Value Features")
    axes[0].set_ylabel("Missing Fraction")
    axes[0].tick_params(axis="x", labelrotation=75)

    feature_var = num_df.var(numeric_only=True).sort_values(ascending=False).head(20)
    feature_var.plot(kind="bar", ax=axes[1], color="#0072B2")
    axes[1].set_title("Top Variance Features")
    axes[1].set_ylabel("Variance")
    axes[1].tick_params(axis="x", labelrotation=75)

    fig.tight_layout()
    fig.savefig(out_dir / "feature_quality_stats.png", bbox_inches="tight")
    plt.close(fig)

    corr_cols = feature_var.index[:12]
    corr = num_df[corr_cols].corr()
    fig, ax = plt.subplots(figsize=(8.5, 7.2))
    sns.heatmap(corr, cmap="coolwarm", center=0, ax=ax)
    ax.set_title("Correlation Heatmap (Top-Variance Features)")
    fig.tight_layout()
    fig.savefig(out_dir / "feature_correlation_heatmap.png", bbox_inches="tight")
    plt.close(fig)


def plot_pca_embeddings(df: pd.DataFrame, all_cols: List[str], out_dir: Path) -> None:
    X = df[all_cols]
    X_imp = SimpleImputer(strategy="median").fit_transform(X)
    X_scaled = StandardScaler().fit_transform(X_imp)
    pca = PCA(n_components=2, random_state=RANDOM_STATE)
    comp = pca.fit_transform(X_scaled)

    proj = pd.DataFrame({"PC1": comp[:, 0], "PC2": comp[:, 1]}, index=df.index)
    proj["dendrite_type"] = df["dendrite_type"].values
    proj["species"] = df["species"].values
    proj["brain_area"] = df["brain_area"].values

    fig, axes = plt.subplots(1, 3, figsize=(16, 4.6))
    sns.scatterplot(data=proj, x="PC1", y="PC2", hue="dendrite_type", s=25, alpha=0.8, ax=axes[0])
    axes[0].set_title("PCA by Dendrite Type")

    sns.scatterplot(data=proj, x="PC1", y="PC2", hue="species", s=25, alpha=0.8, ax=axes[1], palette="Set2")
    axes[1].set_title("PCA by Species")

    top_areas = proj["brain_area"].value_counts().head(8).index
    proj_area = proj[proj["brain_area"].isin(top_areas)]
    sns.scatterplot(data=proj_area, x="PC1", y="PC2", hue="brain_area", s=20, alpha=0.85, ax=axes[2], legend=True)
    axes[2].set_title("PCA by Brain Area (Top 8)")

    for ax in axes:
        ax.legend(loc="best", fontsize=7, frameon=False)

    fig.suptitle(
        f"PCA Embeddings (explained variance: PC1={pca.explained_variance_ratio_[0]:.2f}, "
        f"PC2={pca.explained_variance_ratio_[1]:.2f})",
        fontsize=12,
        fontweight="bold",
    )
    fig.tight_layout()
    fig.savefig(out_dir / "pca_embeddings.png", bbox_inches="tight")
    plt.close(fig)


def plot_performance_grid(
    sweep_results: Dict[Tuple[str, str], pd.DataFrame],
    tasks: List[TaskSpec],
    features: List[FeatureSpec],
    out_dir: Path,
) -> None:
    fig, axes = plt.subplots(len(tasks), len(features), figsize=(16, 13), sharex=True, sharey=True)

    color_map = {"DecisionTree": "#1f77b4", "RandomForest": "#d62728"}
    for r, task in enumerate(tasks):
        for c, feature in enumerate(features):
            ax = axes[r, c]
            res = sweep_results[(task.task_id, feature.feature_id)]
            chance = 1.0 / res["n_classes"].iloc[0]

            for model_name in ["DecisionTree", "RandomForest"]:
                rr = res[res["model"] == model_name]
                ax.plot(rr["node"], rr["train_mean"], color=color_map[model_name], linestyle="-", linewidth=1.5)
                ax.plot(rr["node"], rr["test_mean"], color=color_map[model_name], linestyle="--", linewidth=1.5)
                ax.fill_between(
                    rr["node"],
                    rr["test_mean"] - rr["test_sem"],
                    rr["test_mean"] + rr["test_sem"],
                    color=color_map[model_name],
                    alpha=0.14,
                )

            ax.axhline(chance, linestyle=":", color="gray", linewidth=1.2)
            if r == 0:
                ax.set_title(feature.feature_name, fontweight="bold")
            if c == 0:
                ax.set_ylabel(f"{task.label_name}\nAccuracy")
            if r == len(tasks) - 1:
                ax.set_xlabel("Max Leaf Nodes")
            ax.set_ylim(0.0, 1.02)

    handles = [
        plt.Line2D([], [], color="#1f77b4", linestyle="-", label="DecisionTree train"),
        plt.Line2D([], [], color="#1f77b4", linestyle="--", label="DecisionTree test"),
        plt.Line2D([], [], color="#d62728", linestyle="-", label="RandomForest train"),
        plt.Line2D([], [], color="#d62728", linestyle="--", label="RandomForest test"),
        plt.Line2D([], [], color="gray", linestyle=":", label="Chance"),
    ]
    fig.legend(handles=handles, loc="upper center", ncol=5, bbox_to_anchor=(0.5, 0.99), frameon=False)
    fig.suptitle("Decision Tree and Random Forest Performance Across Tasks", fontsize=16, fontweight="bold", y=1.02)
    fig.tight_layout()
    fig.savefig(out_dir / "decision_tree_performance_grid.png", bbox_inches="tight")
    plt.close(fig)


def fit_best_model_and_plots(
    df: pd.DataFrame,
    feature_spec: FeatureSpec,
    task: TaskSpec,
    sweep_df: pd.DataFrame,
    out_dir: Path,
) -> Dict[str, float]:
    y_raw = df[task.label_col].astype(str)
    y = top_classes(y_raw, task.top_n)
    use_idx = y.notna()
    X = df.loc[use_idx, feature_spec.columns]
    y = y.loc[use_idx]

    best_row = sweep_df.sort_values("test_mean", ascending=False).iloc[0]
    best_model_name = str(best_row["model"])
    best_nodes = int(best_row["node"])

    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=RANDOM_STATE, stratify=y
    )

    model = build_pipeline(best_model_name, best_nodes)
    model.fit(X_train, y_train)
    y_pred = model.predict(X_test)
    acc = accuracy_score(y_test, y_pred)

    labels = sorted(y.unique())
    cm = confusion_matrix(y_test, y_pred, labels=labels, normalize=None)

    fig, ax = plt.subplots(figsize=(6.2, 5.8))
    disp = ConfusionMatrixDisplay(confusion_matrix=cm, display_labels=labels)
    disp.plot(ax=ax, cmap="Blues", colorbar=False, xticks_rotation=45)
    ax.set_title(f"{task.label_name} | {feature_spec.feature_name}\n{best_model_name} ({best_nodes} nodes) | acc={acc:.3f}")
    fig.tight_layout()
    cm_name = f"confusion_{task.task_id}_{feature_spec.feature_id}.png"
    fig.savefig(out_dir / cm_name, bbox_inches="tight")
    plt.close(fig)

    # Feature importance for random forest
    if best_model_name == "RandomForest":
        rf = model.named_steps["model"]
        importances = pd.Series(rf.feature_importances_, index=feature_spec.columns).sort_values(ascending=False).head(20)
        fig, ax = plt.subplots(figsize=(7.0, 5.8))
        importances.sort_values().plot(kind="barh", ax=ax, color="#6a3d9a")
        ax.set_title(f"Top Feature Importances\n{task.label_name} | {feature_spec.feature_name}")
        ax.set_xlabel("Importance")
        fig.tight_layout()
        fi_name = f"feature_importance_{task.task_id}_{feature_spec.feature_id}.png"
        fig.savefig(out_dir / fi_name, bbox_inches="tight")
        plt.close(fig)

    # ROC + PR for binary tasks
    if len(labels) == 2 and hasattr(model.named_steps["model"], "predict_proba"):
        probs = model.predict_proba(X_test)[:, 1]
        y_bin = (y_test == labels[1]).astype(int).to_numpy()
        roc_auc = roc_auc_score(y_bin, probs)
        ap = average_precision_score(y_bin, probs)

        fig, axes = plt.subplots(1, 2, figsize=(11.2, 4.5))
        RocCurveDisplay.from_predictions(y_bin, probs, ax=axes[0], name=f"AUC={roc_auc:.3f}")
        axes[0].set_title(f"ROC | {task.label_name} | {feature_spec.feature_name}")
        PrecisionRecallDisplay.from_predictions(y_bin, probs, ax=axes[1], name=f"AP={ap:.3f}")
        axes[1].set_title(f"Precision-Recall | {task.label_name} | {feature_spec.feature_name}")
        fig.tight_layout()
        roc_pr_name = f"roc_pr_{task.task_id}_{feature_spec.feature_id}.png"
        fig.savefig(out_dir / roc_pr_name, bbox_inches="tight")
        plt.close(fig)

    return {
        "task_id": task.task_id,
        "task_name": task.label_name,
        "feature_id": feature_spec.feature_id,
        "feature_name": feature_spec.feature_name,
        "best_model": best_model_name,
        "best_nodes": best_nodes,
        "cv_test_mean": float(best_row["test_mean"]),
        "cv_test_sem": float(best_row["test_sem"]),
        "test_accuracy": float(acc),
        "n_classes": int(len(labels)),
        "n_samples": int(len(y)),
    }


def create_confusion_montage(tasks: List[TaskSpec], features: List[FeatureSpec], fig_dir: Path) -> None:
    rows = len(tasks)
    cols = len(features)
    fig, axes = plt.subplots(rows, cols, figsize=(16, 13))
    for r, task in enumerate(tasks):
        for c, feature in enumerate(features):
            path = fig_dir / f"confusion_{task.task_id}_{feature.feature_id}.png"
            ax = axes[r, c]
            if path.exists():
                img = plt.imread(path)
                ax.imshow(img)
            ax.axis("off")
            if r == 0:
                ax.set_title(feature.feature_name, fontweight="bold")
            if c == 0:
                ax.set_ylabel(task.label_name, fontweight="bold")
    fig.suptitle("Confusion Matrix Gallery", fontsize=16, fontweight="bold")
    fig.tight_layout()
    fig.savefig(fig_dir / "confusion_matrix_grid.png", bbox_inches="tight")
    plt.close(fig)


def download_thumbnail(url_path: str, out_path: Path) -> bool:
    if not isinstance(url_path, str) or not url_path.startswith("/"):
        return False
    url = f"{CELLTYPES_BASE_URL}{url_path}"
    try:
        response = requests.get(url, timeout=20)
        response.raise_for_status()
        out_path.write_bytes(response.content)
        return True
    except Exception:
        return False


def generate_contact_sheets(df: pd.DataFrame, out_dir: Path) -> None:
    thumbs_dir = out_dir / "thumbs"
    thumbs_dir.mkdir(parents=True, exist_ok=True)

    selected = []
    for dendrite in ["spiny", "aspiny"]:
        part = df[df["dendrite_type"] == dendrite].sample(n=min(8, len(df[df["dendrite_type"] == dendrite])), random_state=RANDOM_STATE)
        selected.append(part)
    selected_df = pd.concat(selected, axis=0).drop_duplicates(subset=["specimen_id"]).head(16)

    morph_files = []
    ephys_files = []
    for _, row in selected_df.iterrows():
        sid = int(row["specimen_id"])
        morph_out = thumbs_dir / f"morph_{sid}.png"
        ephys_out = thumbs_dir / f"ephys_{sid}.png"
        if download_thumbnail(str(row["morph_thumb_path"]), morph_out):
            morph_files.append((morph_out, row["dendrite_type"], row["species"]))
        if download_thumbnail(str(row["ephys_thumb_path"]), ephys_out):
            ephys_files.append((ephys_out, row["dendrite_type"], row["species"]))

    def make_sheet(files: List[Tuple[Path, str, str]], out_name: str, title: str) -> None:
        if not files:
            return
        thumbs = []
        for file_path, dendrite, species in files:
            img = Image.open(file_path).convert("RGB")
            img = ImageOps.pad(img, (320, 210), color=(245, 245, 245))
            thumbs.append((img, dendrite, species))

        cols = 4
        rows = int(np.ceil(len(thumbs) / cols))
        cell_w, cell_h = 360, 250
        canvas = Image.new("RGB", (cols * cell_w, rows * cell_h + 70), (255, 255, 255))

        for i, (img, dendrite, species) in enumerate(thumbs):
            r, c = divmod(i, cols)
            x = c * cell_w + 20
            y = r * cell_h + 45
            canvas.paste(img, (x, y))
            label_bar = Image.new("RGB", (320, 30), (240, 240, 240))
            canvas.paste(label_bar, (x, y + 210))
            draw = ImageOps.expand(Image.new("RGB", (1, 1)), 0)  # dummy to access context consistently
            del draw
            # PIL default font is sufficient for contact labels.
            from PIL import ImageDraw

            text_draw = ImageDraw.Draw(canvas)
            text_draw.text((x + 4, y + 216), f"{dendrite} | {species}", fill=(30, 30, 30))

        from PIL import ImageDraw

        d = ImageDraw.Draw(canvas)
        d.text((16, 14), title, fill=(0, 0, 0))
        canvas.save(out_dir / out_name)

    make_sheet(
        morph_files,
        "morphology_thumbnail_sheet.png",
        "Representative Morphology Thumbnails (Allen Cell Types)",
    )
    make_sheet(
        ephys_files,
        "ephys_thumbnail_sheet.png",
        "Representative Electrophysiology Thumbnails (Allen Cell Types)",
    )


def create_figure_index(fig_dir: Path, summary_df: pd.DataFrame, out_path: Path) -> None:
    cols = list(summary_df.columns)

    def _fmt(val):
        if isinstance(val, float):
            return f"{val:.4f}"
        return str(val)

    header = "| " + " | ".join(cols) + " |"
    divider = "| " + " | ".join(["---"] * len(cols)) + " |"
    rows = [header, divider]
    for _, row in summary_df.iterrows():
        rows.append("| " + " | ".join(_fmt(row[c]) for c in cols) + " |")

    lines = [
        "# ISEF Poster Figure Index",
        "",
        "This file lists generated figures and model-summary highlights.",
        "",
        "## Model Summary",
        "",
    ]
    lines.extend(rows)
    lines.extend(["", "## Figure Files", ""])
    for p in sorted(fig_dir.glob("*.png")):
        lines.append(f"- `{p.name}`")
    out_path.write_text("\n".join(lines))


def add_textbox(slide, left, top, width, height, text, size=20, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Arial"
    return box


def add_bullets(slide, left, top, width, height, bullets, title, title_size=24, body_size=16):
    add_textbox(slide, left, top, width, Inches(0.5), title, size=title_size, bold=True)
    box = slide.shapes.add_textbox(left, top + Inches(0.55), width, height - Inches(0.55))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(body_size)
        p.font.name = "Arial"


def build_poster_draft(fig_dir: Path, out_path: Path, summary_df: pd.DataFrame) -> None:
    prs = Presentation()
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)

    blank = prs.slide_layouts[6]

    # Slide 1: Intro / Methods
    s1 = prs.slides.add_slide(blank)
    add_textbox(
        s1,
        Inches(1.0),
        Inches(0.5),
        Inches(46),
        Inches(1.0),
        "Cell Type Classification Using Electrophysiology and Morphology Features",
        size=34,
        bold=True,
    )
    add_textbox(s1, Inches(1.1), Inches(1.6), Inches(40), Inches(0.6), "ISEF Poster Draft (Template-Matched Layout)", size=18)

    add_bullets(
        s1,
        Inches(1.1),
        Inches(3.0),
        Inches(14.5),
        Inches(12.5),
        [
            "Can neuron cell type be predicted from electrophysiology and morphology-derived features?",
            "Combining feature families should improve generalization compared with using a single modality.",
            "Tree-based models can provide interpretable performance trends as model complexity changes.",
        ],
        title="Hypothesis",
    )
    add_bullets(
        s1,
        Inches(16.7),
        Inches(3.0),
        Inches(14.5),
        Inches(12.5),
        [
            "Data source: Allen Cell Types cached dataset (2,213 filtered cells: spiny/aspiny).",
            "Models: Decision Tree and Random Forest, with max leaf nodes swept from 2 to 15.",
            "Validation: 4-fold stratified cross-validation with train/test accuracy and SEM.",
            "Tasks: Dendrite type, brain area, and merged species+dendrite classes.",
        ],
        title="Methods",
    )

    add_bullets(
        s1,
        Inches(32.3),
        Inches(3.0),
        Inches(14.2),
        Inches(12.5),
        [
            "Allen Institute for Brain Science Cell Types Database",
            "scikit-learn (DecisionTreeClassifier, RandomForestClassifier, StratifiedKFold)",
            "Python scientific stack (NumPy, pandas, matplotlib, seaborn)",
        ],
        title="References",
    )

    for img_name, left, top, width in [
        ("dataset_overview.png", 1.1, 17.0, 14.5),
        ("morphology_thumbnail_sheet.png", 16.7, 17.0, 14.5),
        ("pca_embeddings.png", 32.3, 17.0, 14.2),
    ]:
        p = fig_dir / img_name
        if p.exists():
            s1.shapes.add_picture(str(p), Inches(left), Inches(top), width=Inches(width))

    # Slide 2: Results
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, Inches(1.0), Inches(0.5), Inches(46), Inches(1.0), "Results", size=34, bold=True)

    top_line = summary_df.sort_values("cv_test_mean", ascending=False).head(3)
    bullets = []
    for _, row in top_line.iterrows():
        bullets.append(
            f"{row['task_name']} + {row['feature_name']}: {row['best_model']} ({int(row['best_nodes'])} nodes), "
            f"CV test={row['cv_test_mean']:.3f}±{row['cv_test_sem']:.3f}"
        )
    add_bullets(
        s2,
        Inches(1.1),
        Inches(2.0),
        Inches(45.0),
        Inches(2.4),
        bullets,
        title="Key Findings",
        title_size=22,
        body_size=15,
    )

    grid = fig_dir / "decision_tree_performance_grid.png"
    if grid.exists():
        s2.shapes.add_picture(str(grid), Inches(1.1), Inches(5.0), width=Inches(29.8))

    cm_grid = fig_dir / "confusion_matrix_grid.png"
    if cm_grid.exists():
        s2.shapes.add_picture(str(cm_grid), Inches(31.4), Inches(5.0), width=Inches(15.2))

    roc_pr = fig_dir / "roc_pr_dendrite_type_electrophysiology.png"
    if roc_pr.exists():
        s2.shapes.add_picture(str(roc_pr), Inches(31.4), Inches(22.0), width=Inches(15.2))

    fi = fig_dir / "feature_importance_dendrite_type_all_features.png"
    if fi.exists():
        s2.shapes.add_picture(str(fi), Inches(1.1), Inches(26.2), width=Inches(20.5))

    ephys_sheet = fig_dir / "ephys_thumbnail_sheet.png"
    if ephys_sheet.exists():
        s2.shapes.add_picture(str(ephys_sheet), Inches(22.0), Inches(26.2), width=Inches(24.6))

    prs.save(str(out_path))


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate figure bank + poster draft for ISEF.")
    parser.add_argument("--data-dir", type=Path, default=PROJECT_ROOT / "scripts/cell_types")
    parser.add_argument("--out-dir", type=Path, default=PROJECT_ROOT / "poster/output")
    args = parser.parse_args()

    out_dir: Path = args.out_dir
    fig_dir = out_dir / "figures"
    fig_dir.mkdir(parents=True, exist_ok=True)

    df, ephys_cols, morph_cols, all_cols = load_local_data(args.data_dir)
    plot_dataset_overview(df, fig_dir)
    plot_feature_statistics(df, ephys_cols, fig_dir)
    plot_pca_embeddings(df, all_cols, fig_dir)
    generate_contact_sheets(df, fig_dir)

    tasks = [
        TaskSpec(task_id="dendrite_type", label_col="dendrite_type", label_name="Dendrite Type"),
        TaskSpec(task_id="brain_area", label_col="brain_area", label_name="Brain Area", top_n=8),
        TaskSpec(
            task_id="species_dendrite",
            label_col="species_dendrite",
            label_name="Merged Species + Dendrite",
        ),
    ]
    features = [
        FeatureSpec(feature_id="electrophysiology", feature_name="Electrophysiology", columns=ephys_cols),
        FeatureSpec(feature_id="morphology", feature_name="Morphology", columns=morph_cols),
        FeatureSpec(feature_id="all_features", feature_name="All Features", columns=all_cols),
    ]
    node_values = np.arange(2, 16)

    sweep_results: Dict[Tuple[str, str], pd.DataFrame] = {}
    summary_rows: List[Dict[str, float]] = []

    for task in tasks:
        y_full = top_classes(df[task.label_col].astype(str), task.top_n)
        for feature in features:
            X = df[feature.columns]
            y = y_full.copy()
            use_idx = y.notna()
            X = X.loc[use_idx]
            y = y.loc[use_idx]

            sweep_df = evaluate_node_sweep(X, y, node_values=node_values, n_splits=4)
            sweep_df["task_id"] = task.task_id
            sweep_df["feature_id"] = feature.feature_id
            sweep_df["n_classes"] = int(y.nunique())
            sweep_results[(task.task_id, feature.feature_id)] = sweep_df

            sweep_file = fig_dir / f"node_sweep_{task.task_id}_{feature.feature_id}.csv"
            sweep_df.to_csv(sweep_file, index=False)

            fig, ax = plt.subplots(figsize=(7.0, 5.0))
            for model_name, color in [("DecisionTree", "#1f77b4"), ("RandomForest", "#d62728")]:
                rr = sweep_df[sweep_df["model"] == model_name]
                ax.plot(rr["node"], rr["train_mean"], color=color, linestyle="-", label=f"{model_name} train")
                ax.plot(rr["node"], rr["test_mean"], color=color, linestyle="--", label=f"{model_name} test")
                ax.fill_between(rr["node"], rr["test_mean"] - rr["test_sem"], rr["test_mean"] + rr["test_sem"], color=color, alpha=0.15)
            chance = 1.0 / y.nunique()
            ax.axhline(chance, color="gray", linestyle=":", label="Chance")
            ax.set_xlabel("Max Leaf Nodes")
            ax.set_ylabel("Accuracy")
            ax.set_ylim(0.0, 1.02)
            ax.set_title(f"{task.label_name} | {feature.feature_name}")
            ax.legend(fontsize=8, frameon=False)
            fig.tight_layout()
            fig.savefig(fig_dir / f"node_sweep_{task.task_id}_{feature.feature_id}.png", bbox_inches="tight")
            plt.close(fig)

            row = fit_best_model_and_plots(df, feature, task, sweep_df, fig_dir)
            summary_rows.append(row)

    plot_performance_grid(sweep_results, tasks, features, fig_dir)
    create_confusion_montage(tasks, features, fig_dir)

    summary_df = pd.DataFrame(summary_rows).sort_values(["task_id", "feature_id"]).reset_index(drop=True)
    summary_df.to_csv(out_dir / "model_summary.csv", index=False)

    create_figure_index(fig_dir, summary_df, out_dir / "FIGURE_INDEX.md")
    build_poster_draft(fig_dir, out_dir / "ISEF_Poster_Draft.pptx", summary_df)

    print(f"Generated assets in: {out_dir.resolve()}")
    print(f"Figure count: {len(list(fig_dir.glob('*.png')))}")
    print(f"Poster draft: {(out_dir / 'ISEF_Poster_Draft.pptx').resolve()}")


if __name__ == "__main__":
    main()
